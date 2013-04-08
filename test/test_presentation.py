# -*- coding: utf-8 -*-
#
# test_presentation.py
#
# Copyright (C) 2012, 2013 Steve Canny scanny@cisco.com
#
# This module is part of python-pptx and is released under
# the MIT License: http://www.opensource.org/licenses/mit-license.php

"""Test suite for pptx.presentation module."""

import gc
import os

from hamcrest import assert_that, is_, is_in, is_not, equal_to
from StringIO import StringIO

from mock import Mock

import pptx.presentation

from pptx.exceptions import InvalidPackageError
from pptx.oxml import oxml_fromstring, oxml_tostring, oxml_parse
from pptx.packaging import prettify_nsdecls
from pptx.presentation import (
    Package, _RelationshipCollection, _Relationship, Presentation,
    PartCollection, BasePart, Part, SlideCollection, BaseSlide, Slide,
    SlideLayout, SlideMaster, Image)
from pptx.shapes import ShapeCollection
from pptx.spec import namespaces, qtag
from pptx.spec import (
    CT_PRESENTATION, CT_SLIDE, CT_SLIDE_LAYOUT, CT_SLIDE_MASTER)
from pptx.spec import (
    RT_IMAGE, RT_OFFICE_DOCUMENT, RT_PRES_PROPS, RT_SLIDE, RT_SLIDE_LAYOUT,
    RT_SLIDE_MASTER)
from testing import TestCase

import logging
log = logging.getLogger('pptx.test.presentation')
log.setLevel(logging.DEBUG)
# log.setLevel(logging.INFO)
ch = logging.StreamHandler()
ch.setLevel(logging.DEBUG)
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - '
                              '%(message)s')
ch.setFormatter(formatter)
log.addHandler(ch)


# module globals -------------------------------------------------------------
def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
test_file_dir = absjoin(thisdir, 'test_files')

test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_bmp_path = absjoin(test_file_dir, 'python.bmp')
new_image_path = absjoin(test_file_dir, 'monty-truth.png')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')
images_pptx_path = absjoin(test_file_dir, 'with_images.pptx')

nsmap = namespaces('a', 'r', 'p')
nsprefix_decls = (
    ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xm'
    'lns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="h'
    'ttp://schemas.openxmlformats.org/officeDocument/2006/relationships"')


def _sldLayout1():
    path = os.path.join(thisdir, 'test_files/slideLayout1.xml')
    sldLayout = oxml_parse(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = ShapeCollection(spTree)
    return shapes


class PartBuilder(object):
    """Builder class for test Parts"""
    def __init__(self):
        self.partname = '/ppt/slides/slide1.xml'

    def with_partname(self, partname):
        self.partname = partname
        return self

    def build(self):
        p = BasePart()
        p.partname = self.partname
        return p


class RelationshipCollectionBuilder(object):
    """Builder class for test RelationshipCollections"""
    partname_tmpls = {RT_SLIDE_MASTER: '/ppt/slideMasters/slideMaster%d.xml',
                      RT_SLIDE: '/ppt/slides/slide%d.xml'}

    def __init__(self):
        self.relationships = []
        self.next_rel_num = 1
        self.next_partnums = {}
        self.reltype_ordering = None

    def with_ordering(self, *reltypes):
        self.reltype_ordering = tuple(reltypes)
        return self

    def with_tuple_targets(self, count, reltype):
        for i in range(count):
            rId = self.__next_rId
            partname = self.__next_tuple_partname(reltype)
            target = PartBuilder().with_partname(partname).build()
            rel = _Relationship(rId, reltype, target)
            self.relationships.append(rel)
        return self

    # def with_singleton_target(self, reltype):
    #     rId = self.__next_rId
    #     partname = self.__singleton_partname(reltype)
    #     target = PartBuilder().with_partname(partname).build()
    #     rel = _Relationship(rId, reltype, target)
    #     self.relationships.append(rel)
    #     return self
    #
    def __next_partnum(self, reltype):
        if reltype not in self.next_partnums:
            self.next_partnums[reltype] = 1
        partnum = self.next_partnums[reltype]
        self.next_partnums[reltype] = partnum + 1
        return partnum

    @property
    def __next_rId(self):
        rId = 'rId%d' % self.next_rel_num
        self.next_rel_num += 1
        return rId

    def __next_tuple_partname(self, reltype):
        partname_tmpl = self.partname_tmpls[reltype]
        partnum = self.__next_partnum(reltype)
        return partname_tmpl % partnum

    def build(self):
        rels = _RelationshipCollection()
        for rel in self.relationships:
            rels._additem(rel)
        if self.reltype_ordering:
            rels._reltype_ordering = self.reltype_ordering
        return rels


class TestBasePart(TestCase):
    """Test BasePart"""
    def setUp(self):
        self.basepart = BasePart()
        self.cls = BasePart

    def test__add_relationship_adds_specified_relationship(self):
        """BasePart._add_relationship adds specified relationship"""
        # setup -----------------------
        reltype = RT_IMAGE
        target = Mock(name='image')
        # exercise --------------------
        rel = self.basepart._add_relationship(reltype, target)
        # verify ----------------------
        expected = ('rId1', reltype, target)
        actual = (rel._rId, rel._reltype, rel._target)
        msg = "\nExpected: %s\n     Got: %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__blob_value_for_binary_part(self):
        """BasePart._blob value is correct for binary part"""
        # setup -----------------------
        blob = '0123456789'
        self.basepart._load_blob = blob
        self.basepart.partname = '/docProps/thumbnail.jpeg'
        # exercise --------------------
        retval = self.basepart._blob
        # verify ----------------------
        expected = blob
        actual = retval
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__blob_value_for_xml_part(self):
        """BasePart._blob value is correct for XML part"""
        # setup -----------------------
        elm = oxml_fromstring('<root><elm1 attr="one"/></root>')
        self.basepart._element = elm
        self.basepart.partname = '/ppt/presentation.xml'
        # exercise --------------------
        retval = self.basepart._blob
        # verify ----------------------
        expected = "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>"\
                   '\n<root>\n  <elm1 attr="one"/>\n</root>\n'
        actual = retval
        msg = "expected: \n'%s'\n, got \n'%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__content_type_raises_on_accessed_before_assigned(self):
        """BasePart._content_type raises on access before assigned"""
        with self.assertRaises(ValueError):
            self.basepart._content_type

    def test__load_sets__element_for_xml_part(self):
        """BasePart._load() sets _element for xml part"""
        # setup -----------------------
        pkgpart = Mock(name='pptx.packaging.Part')
        pkgpart.partname = '/ppt/presentation.xml'
        pkgpart.blob = '<root><elm1   attr="spam"/></root>'
        pkgpart.relationships = []
        part_dict = {}
        part = self.basepart._load(pkgpart, part_dict)
        # exercise --------------------
        elm = part._element
        # verify ----------------------
        expected = '<root><elm1 attr="spam"/></root>'
        actual = oxml_tostring(elm)
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_observable_on_partname(self):
        """BasePart observable on partname value change"""
        # setup -----------------------
        old_partname = '/ppt/slides/slide1.xml'
        new_partname = '/ppt/slides/slide2.xml'
        observer = Mock()
        self.basepart.partname = old_partname
        self.basepart.add_observer(observer)
        # exercise --------------------
        self.basepart.partname = new_partname
        # verify ----------------------
        observer.notify.assert_called_with(self.basepart, 'partname',
                                           new_partname)

    def test_partname_setter(self):
        """BasePart.partname setter stores passed value"""
        # setup -----------------------
        partname = '/ppt/presentation.xml'
        # exercise ----------------
        self.basepart.partname = partname
        # verify ------------------
        expected = partname
        actual = self.basepart.partname
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class TestBaseSlide(TestCase):
    """Test BaseSlide"""
    def setUp(self):
        self.base_slide = BaseSlide()

    def test_name_value(self):
        """BaseSlide.name value is correct"""
        # setup -----------------------
        self.base_slide._element = _sldLayout1()
        # exercise --------------------
        name = self.base_slide.name
        # verify ----------------------
        expected = 'Title Slide'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_shapes_size_after__load(self):
        """BaseSlide.shapes is expected size after _load()"""
        # setup -----------------------
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        pkgpart = Mock(name='pptx.packaging.Part')
        pkgpart.partname = '/ppt/slides/slide1.xml'
        with open(path, 'r') as f:
            pkgpart.blob = f.read()
        pkgpart.relationships = []
        part_dict = {}
        self.base_slide._load(pkgpart, part_dict)
        # exercise --------------------
        shapes = self.base_slide.shapes
        # verify ----------------------
        self.assertLength(shapes, 9)


class TestImage(TestCase):
    """Test Image"""
    def test_construction_from_file(self):
        """Image(path) constructor produces correct attribute values"""
        # exercise --------------------
        image = Image(test_image_path)
        # verify ----------------------
        assert_that(image.ext, is_(equal_to('.jpeg')))
        assert_that(image._content_type, is_(equal_to('image/jpeg')))
        assert_that(len(image._blob), is_(equal_to(3277)))

    def test_construction_from_stream(self):
        """Image(stream) construction produces correct attribute values"""
        # exercise --------------------
        with open(test_image_path) as f:
            stream = StringIO(f.read())
        image = Image(stream)
        # verify ----------------------
        assert_that(image.ext, is_(equal_to('.jpg')))
        assert_that(image._content_type, is_(equal_to('image/jpeg')))
        assert_that(len(image._blob), is_(equal_to(3277)))

    def test_construction_from_file_raises_on_bad_path(self):
        """Image(path) constructor raises on bad path"""
        # verify ----------------------
        with self.assertRaises(IOError):
            Image('foobar27.png')

    def test___ext_from_image_stream_raises_on_incompatible_format(self):
        """Image.__ext_from_image_stream() raises on incompatible format"""
        # verify ----------------------
        with self.assertRaises(ValueError):
            with open(test_bmp_path) as stream:
                Image._Image__ext_from_image_stream(stream)

    def test___image_ext_content_type_known_type(self):
        """Image.__image_ext_content_type() correct for known content type"""
        # exercise --------------------
        content_type = Image._Image__image_ext_content_type('.jpeg')
        # verify ----------------------
        expected = 'image/jpeg'
        actual = content_type
        msg = ("expected content type '%s', got '%s'" % (expected, actual))
        self.assertEqual(expected, actual, msg)

    def test___image_ext_content_type_raises_on_bad_ext(self):
        """Image.__image_ext_content_type() raises on bad extension"""
        # verify ----------------------
        with self.assertRaises(TypeError):
            Image._Image__image_ext_content_type('.xj7')

    def test___image_ext_content_type_raises_on_non_img_ext(self):
        """Image.__image_ext_content_type() raises on non-image extension"""
        # verify ----------------------
        with self.assertRaises(TypeError):
            Image._Image__image_ext_content_type('.xml')


class TestImageCollection(TestCase):
    """Test ImageCollection"""
    def test_add_image_returns_matching_image(self):
        """ImageCollection.add_image() returns existing image on match"""
        # setup -----------------------
        pkg = Package(images_pptx_path)
        matching_idx = 4
        matching_image = pkg._images[matching_idx]
        # exercise --------------------
        image = pkg._images.add_image(test_image_path)
        # verify ----------------------
        expected = matching_image
        actual = image
        msg = ("expected images[%d], got images[%d]"
               % (matching_idx, pkg._images.index(image)))
        self.assertEqual(expected, actual, msg)

    def test_add_image_adds_new_image(self):
        """ImageCollection.add_image() adds new image on no match"""
        # setup -----------------------
        pkg = Package(images_pptx_path)
        expected_partname = '/ppt/media/image8.png'
        expected_len = len(pkg._images) + 1
        expected_sha1 = '79769f1e202add2e963158b532e36c2c0f76a70c'
        # exercise --------------------
        image = pkg._images.add_image(new_image_path)
        # verify ----------------------
        expected = (expected_partname, expected_len, expected_sha1)
        actual = (image.partname, len(pkg._images), image._sha1)
        msg = "\nExpected: %s\n     Got: %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class TestPackage(TestCase):
    """Test Package"""
    def setUp(self):
        self.test_pptx_path = absjoin(test_file_dir, 'test_python-pptx.pptx')
        if os.path.isfile(self.test_pptx_path):
            os.remove(self.test_pptx_path)

    def tearDown(self):
        if os.path.isfile(self.test_pptx_path):
            os.remove(self.test_pptx_path)

    def test_construction_with_no_path_loads_default_template(self):
        """Package() call with no path loads default template"""
        prs = Package().presentation
        assert_that(prs, is_not(None))
        slidemasters = prs.slidemasters
        assert_that(slidemasters, is_not(None))
        assert_that(len(slidemasters), is_(1))
        slidelayouts = slidemasters[0].slidelayouts
        assert_that(slidelayouts, is_not(None))
        assert_that(len(slidelayouts), is_(11))

    def test_instances_are_tracked(self):
        """Package instances are tracked"""
        pkg = Package()
        self.assertIn(pkg, Package.instances())

    def test_instance_refs_are_garbage_collected(self):
        """Package instances are tracked"""
        pkg = Package()
        pkg1_repr = "%r" % pkg
        pkg = Package()
        # pkg2_repr = "%r" % pkg
        gc.collect()
        reprs = [repr(pkg_inst) for pkg_inst in Package.instances()]
        # log.debug("pkg1, pkg2, reprs: %s, %s, %s"
        #           % (pkg1_repr, pkg2_repr, reprs))
        assert_that(pkg1_repr, is_not(is_in(reprs)))

    def test_containing_returns_correct_pkg(self):
        """Package.containing() returns right package instance"""
        # setup -----------------------
        pkg1 = Package(test_pptx_path)
        pkg1.presentation  # does nothing, just needed to fake out pep8 warning
        pkg2 = Package(test_pptx_path)
        slide = pkg2.presentation.slides[0]
        # exercise --------------------
        found_pkg = Package.containing(slide)
        # verify ----------------------
        expected = pkg2
        actual = found_pkg
        msg = "expected %r, got %r" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_containing_raises_on_no_pkg_contains_part(self):
        """Package.containing(part) raises on no package contains part"""
        # setup -----------------------
        pkg = Package(test_pptx_path)
        pkg.presentation  # does nothing, just needed to fake out pep8 warning
        part = Mock(name='part')
        # verify ----------------------
        with self.assertRaises(KeyError):
            Package.containing(part)

    def test_open_gathers_image_parts(self):
        """Package open gathers image parts into image collection"""
        # exercise --------------------
        pkg = Package(images_pptx_path)
        # verify ----------------------
        expected = 7
        actual = len(pkg._Package__images)
        msg = "expected image count of %d, got %d" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_presentation_presentation_after_open(self):
        """Package.presentation is instance of Presentation after open()"""
        # setup -----------------------
        cls = Presentation
        pkg = Package()
        # exercise --------------------
        obj = pkg.presentation
        # verify ----------------------
        actual = isinstance(obj, cls)
        msg = ("expected instance of '%s', got type '%s'"
               % (cls.__name__, type(obj).__name__))
        self.assertTrue(actual, msg)

    def test_saved_file_has_plausible_contents(self):
        """Package.save produces a .pptx with plausible contents"""
        # setup -----------------------
        pkg = Package()
        # exercise --------------------
        pkg.save(self.test_pptx_path)
        # verify ----------------------
        pkg = Package(self.test_pptx_path)
        prs = pkg.presentation
        assert_that(prs, is_not(None))
        slidemasters = prs.slidemasters
        assert_that(slidemasters, is_not(None))
        assert_that(len(slidemasters), is_(1))
        slidelayouts = slidemasters[0].slidelayouts
        assert_that(slidelayouts, is_not(None))
        assert_that(len(slidelayouts), is_(11))


class TestPart(TestCase):
    """Test Part"""
    def test_constructs_presentation_for_rt_officedocument(self):
        """Part() returns Presentation for RT_OFFICE_DOCUMENT"""
        # setup -----------------------
        cls = Presentation
        # exercise --------------------
        obj = Part(RT_OFFICE_DOCUMENT, CT_PRESENTATION)
        # verify ----------------------
        self.assertIsInstance(obj, cls)

    def test_constructs_slide_for_rt_slide(self):
        """Part() returns Slide for RT_SLIDE"""
        # setup -----------------------
        cls = Slide
        # exercise --------------------
        obj = Part(RT_SLIDE, CT_SLIDE)
        # verify ----------------------
        self.assertIsInstance(obj, cls)

    def test_constructs_slidelayout_for_rt_slidelayout(self):
        """Part() returns SlideLayout for RT_SLIDE_LAYOUT"""
        # setup -----------------------
        cls = SlideLayout
        # exercise --------------------
        obj = Part(RT_SLIDE_LAYOUT, CT_SLIDE_LAYOUT)
        # verify ----------------------
        self.assertIsInstance(obj, cls)

    def test_constructs_slidemaster_for_rt_slidemaster(self):
        """Part() returns SlideMaster for RT_SLIDE_MASTER"""
        # setup -----------------------
        cls = SlideMaster
        # exercise --------------------
        obj = Part(RT_SLIDE_MASTER, CT_SLIDE_MASTER)
        # verify ----------------------
        self.assertIsInstance(obj, cls)

    def test_contructor_raises_on_invalid_prs_content_type(self):
        """Part() raises on invalid presentation content type"""
        with self.assertRaises(InvalidPackageError):
            Part(RT_OFFICE_DOCUMENT, CT_SLIDE_MASTER)


class TestPartCollection(TestCase):
    """Test PartCollection"""
    def test__loadpart_sorts_loaded_parts(self):
        """PartCollection._loadpart sorts loaded parts"""
        # setup -----------------------
        partname1 = '/ppt/slides/slide1.xml'
        partname2 = '/ppt/slides/slide2.xml'
        partname3 = '/ppt/slides/slide3.xml'
        part1 = Mock(name='part1')
        part1.partname = partname1
        part2 = Mock(name='part2')
        part2.partname = partname2
        part3 = Mock(name='part3')
        part3.partname = partname3
        parts = PartCollection()
        # exercise --------------------
        parts._loadpart(part2)
        parts._loadpart(part3)
        parts._loadpart(part1)
        # verify ----------------------
        expected = [partname1, partname2, partname3]
        actual = [part.partname for part in parts]
        msg = "expected %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class TestPresentation(TestCase):
    """Test Presentation"""
    def setUp(self):
        self.prs = Presentation()

    def test__blob_rewrites_sldIdLst(self):
        """Presentation._blob rewrites sldIdLst"""
        # setup -----------------------
        rels = RelationshipCollectionBuilder()
        rels = rels.with_tuple_targets(2, RT_SLIDE_MASTER)
        rels = rels.with_tuple_targets(3, RT_SLIDE)
        rels = rels.with_ordering(RT_SLIDE_MASTER, RT_SLIDE)
        rels = rels.build()
        prs = Presentation()
        prs._relationships = rels
        prs.partname = '/ppt/presentation.xml'
        path = os.path.join(thisdir, 'test_files/presentation.xml')
        prs._element = oxml_parse(path).getroot()
        # exercise --------------------
        blob = prs._blob
        # verify ----------------------
        presentation = oxml_fromstring(blob)
        sldIds = presentation.xpath('./p:sldIdLst/p:sldId', namespaces=nsmap)
        expected = ['rId3', 'rId4', 'rId5']
        actual = [sldId.get(qtag('r:id')) for sldId in sldIds]
        msg = "expected ordering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_slidemasters_property_empty_on_construction(self):
        """Presentation.slidemasters property empty on construction"""
        # verify ----------------------
        self.assertIsSizedProperty(self.prs, 'slidemasters', 0)

    def test_slidemasters_correct_length_after_pkg_open(self):
        """Presentation.slidemasters correct length after load"""
        # setup -----------------------
        pkg = Package(test_pptx_path)
        prs = pkg.presentation
        # exercise --------------------
        slidemasters = prs.slidemasters
        # verify ----------------------
        self.assertLength(slidemasters, 1)

    def test_slides_property_empty_on_construction(self):
        """Presentation.slides property empty on construction"""
        # verify ----------------------
        self.assertIsSizedProperty(self.prs, 'slides', 0)

    def test_slides_correct_length_after_pkg_open(self):
        """Presentation.slides correct length after load"""
        # setup -----------------------
        pkg = Package(test_pptx_path)
        prs = pkg.presentation
        # exercise --------------------
        slides = prs.slides
        # verify ----------------------
        self.assertLength(slides, 1)


class Test_Relationship(TestCase):
    """Test _Relationship"""
    def setUp(self):
        rId = 'rId1'
        reltype = RT_SLIDE
        target_part = None
        self.rel = _Relationship(rId, reltype, target_part)

    def test_constructor_raises_on_bad_rId(self):
        """_Relationship constructor raises on non-standard rId"""
        with self.assertRaises(AssertionError):
            _Relationship('Non-std14', None, None)

    def test__num_value(self):
        """_Relationship._num value is correct"""
        # setup -----------------------
        num = 91
        rId = 'rId%d' % num
        rel = _Relationship(rId, None, None)
        # verify ----------------------
        assert_that(rel._num, is_(equal_to(num)))

    def test__num_value_on_non_standard_rId(self):
        """_Relationship._num value is correct for non-standard rId"""
        # setup -----------------------
        rel = _Relationship('rIdSm', None, None)
        # verify ----------------------
        assert_that(rel._num, is_(equal_to(9999)))

    def test__rId_setter(self):
        """Relationship._rId setter stores passed value"""
        # setup -----------------------
        rId = 'rId9'
        # exercise ----------------
        self.rel._rId = rId
        # verify ------------------
        expected = rId
        actual = self.rel._rId
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class Test_RelationshipCollection(TestCase):
    """Test _RelationshipCollection"""
    def setUp(self):
        self.relationships = _RelationshipCollection()

    def __reltype_ordering_mock(self):
        """
        Return RelationshipCollection instance with mocked-up contents
        suitable for testing _reltype_ordering.
        """
        # setup -----------------------
        partnames = ['/ppt/slides/slide4.xml',
                     '/ppt/slideLayouts/slideLayout1.xml',
                     '/ppt/slideMasters/slideMaster1.xml',
                     '/ppt/slides/slide1.xml',
                     '/ppt/presProps.xml']
        part1 = Mock(name='part1')
        part1.partname = partnames[0]
        part2 = Mock(name='part2')
        part2.partname = partnames[1]
        part3 = Mock(name='part3')
        part3.partname = partnames[2]
        part4 = Mock(name='part4')
        part4.partname = partnames[3]
        part5 = Mock(name='part5')
        part5.partname = partnames[4]
        rel1 = _Relationship('rId1', RT_SLIDE,        part1)
        rel2 = _Relationship('rId2', RT_SLIDE_LAYOUT, part2)
        rel3 = _Relationship('rId3', RT_SLIDE_MASTER, part3)
        rel4 = _Relationship('rId4', RT_SLIDE,        part4)
        rel5 = _Relationship('rId5', RT_PRES_PROPS,   part5)
        relationships = _RelationshipCollection()
        relationships._additem(rel1)
        relationships._additem(rel2)
        relationships._additem(rel3)
        relationships._additem(rel4)
        relationships._additem(rel5)
        return (relationships, partnames)

    def test__additem_raises_on_dup_rId(self):
        """_RelationshipCollection._additem raises on duplicate rId"""
        # setup -----------------------
        part1 = BasePart()
        part2 = BasePart()
        rel1 = _Relationship('rId9', None, part1)
        rel2 = _Relationship('rId9', None, part2)
        self.relationships._additem(rel1)
        # verify ----------------------
        with self.assertRaises(ValueError):
            self.relationships._additem(rel2)

    def test__additem_maintains_rId_ordering(self):
        """_RelationshipCollection maintains rId ordering on additem()"""
        # setup -----------------------
        part1 = BasePart()
        part2 = BasePart()
        part3 = BasePart()
        rel1 = _Relationship('rId1', None, part1)
        rel2 = _Relationship('rId2', None, part2)
        rel3 = _Relationship('rId3', None, part3)
        # exercise --------------------
        self.relationships._additem(rel2)
        self.relationships._additem(rel1)
        self.relationships._additem(rel3)
        # verify ----------------------
        expected = ['rId1', 'rId2', 'rId3']
        actual = [rel._rId for rel in self.relationships]
        msg = "expected ordering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__additem_maintains_reltype_ordering(self):
        """_RelationshipCollection maintains reltype ordering on additem()"""
        # setup -----------------------
        relationships, partnames = self.__reltype_ordering_mock()
        ordering = (RT_SLIDE_MASTER, RT_SLIDE_LAYOUT, RT_SLIDE)
        relationships._reltype_ordering = ordering
        partname = '/ppt/slides/slide2.xml'
        part = Mock(name='new_part')
        part.partname = partname
        rId = relationships._next_rId
        rel = _Relationship(rId, RT_SLIDE, part)
        # exercise --------------------
        relationships._additem(rel)
        # verify ordering -------------
        expected = [partnames[2], partnames[1], partnames[3],
                    partname, partnames[0], partnames[4]]
        actual = [r._target.partname for r in relationships]
        msg = "expected ordering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_rels_of_reltype_return_value(self):
        """RelationshipCollection._rels_of_reltype returns correct rels"""
        # setup -----------------------
        relationships, partnames = self.__reltype_ordering_mock()
        # exercise --------------------
        retval = relationships.rels_of_reltype(RT_SLIDE)
        # verify ordering -------------
        expected = ['rId1', 'rId4']
        actual = [rel._rId for rel in retval]
        msg = "expected %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__reltype_ordering_sorts_rels(self):
        """RelationshipCollection._reltype_ordering sorts rels"""
        # setup -----------------------
        relationships, partnames = self.__reltype_ordering_mock()
        ordering = (RT_SLIDE_MASTER, RT_SLIDE_LAYOUT, RT_SLIDE)
        # exercise --------------------
        relationships._reltype_ordering = ordering
        # verify ordering -------------
        assert_that(relationships._reltype_ordering, is_(equal_to(ordering)))
        expected = [partnames[2], partnames[1], partnames[3], partnames[0],
                    partnames[4]]
        actual = [rel._target.partname for rel in relationships]
        msg = "expected ordering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__reltype_ordering_renumbers_rels(self):
        """RelationshipCollection._reltype_ordering renumbers rels"""
        # setup -----------------------
        relationships, partnames = self.__reltype_ordering_mock()
        ordering = (RT_SLIDE_MASTER, RT_SLIDE_LAYOUT, RT_SLIDE)
        # exercise --------------------
        relationships._reltype_ordering = ordering
        # verify renumbering ----------
        expected = ['rId1', 'rId2', 'rId3', 'rId4', 'rId5']
        actual = [rel._rId for rel in relationships]
        msg = "expected numbering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__next_rId_fills_gap(self):
        """_RelationshipCollection._next_rId fills gap in rId sequence"""
        # setup -----------------------
        part1 = BasePart()
        part2 = BasePart()
        part3 = BasePart()
        part4 = BasePart()
        rel1 = _Relationship('rId1', None, part1)
        rel2 = _Relationship('rId2', None, part2)
        rel3 = _Relationship('rId3', None, part3)
        rel4 = _Relationship('rId4', None, part4)
        cases = (('rId1', (rel2, rel3, rel4)),
                 ('rId2', (rel1, rel3, rel4)),
                 ('rId3', (rel1, rel2, rel4)),
                 ('rId4', (rel1, rel2, rel3)))
        # exercise --------------------
        expected_rIds = []
        actual_rIds = []
        for expected_rId, rels in cases:
            expected_rIds.append(expected_rId)
            relationships = _RelationshipCollection()
            for rel in rels:
                relationships._additem(rel)
            actual_rIds.append(relationships._next_rId)
        # verify ----------------------
        expected = expected_rIds
        actual = actual_rIds
        msg = "expected rIds %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_reorders_on_partname_change(self):
        """RelationshipCollection reorders on partname change"""
        # setup -----------------------
        partname1 = '/ppt/slides/slide1.xml'
        partname2 = '/ppt/slides/slide2.xml'
        partname3 = '/ppt/slides/slide3.xml'
        part1 = PartBuilder().with_partname(partname1).build()
        part2 = PartBuilder().with_partname(partname2).build()
        rel1 = _Relationship('rId1', RT_SLIDE, part1)
        rel2 = _Relationship('rId2', RT_SLIDE, part2)
        relationships = _RelationshipCollection()
        relationships._reltype_ordering = (RT_SLIDE)
        relationships._additem(rel1)
        relationships._additem(rel2)
        # exercise --------------------
        part1.partname = partname3
        # verify ----------------------
        expected = [partname2, partname3]
        actual = [rel._target.partname for rel in relationships]
        msg = "expected ordering %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class TestSlide(TestCase):
    """Test Slide"""
    def setUp(self):
        self.sld = Slide()

    def test_constructor_sets_correct_content_type(self):
        """Slide constructor sets correct content type"""
        # exercise --------------------
        content_type = self.sld._content_type
        # verify ----------------------
        expected = CT_SLIDE
        actual = content_type
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_construction_adds_slide_layout_relationship(self):
        """Slide(slidelayout) adds relationship slide->slidelayout"""
        # setup -----------------------
        slidelayout = SlideLayout()
        slidelayout._shapes = _sldLayout1_shapes()
        # exercise --------------------
        slide = Slide(slidelayout)
        # verify length ---------------
        expected = 1
        actual = len(slide._relationships)
        msg = ("expected len(slide._relationships) of %d, got %d"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)
        # verify values ---------------
        rel = slide._relationships[0]
        expected = ('rId1', RT_SLIDE_LAYOUT, slidelayout)
        actual = (rel._rId, rel._reltype, rel._target)
        msg = "expected relationship\n%s\ngot\n%s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__element_minimal_sld_on_construction(self):
        """Slide._element is minimal sld on construction"""
        # setup -----------------------
        path = os.path.join(thisdir, 'test_files/minimal_slide.xml')
        # exercise --------------------
        elm = self.sld._element
        # verify ----------------------
        with open(path, 'r') as f:
            expected = f.read()
        actual = prettify_nsdecls(
            oxml_tostring(elm, encoding='UTF-8', pretty_print=True,
                          standalone=True))
        msg = "\nexpected:\n\n'%s'\n\nbut got:\n\n'%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_slidelayout_property_none_on_construction(self):
        """Slide.slidelayout property None on construction"""
        # verify ----------------------
        self.assertIsProperty(self.sld, 'slidelayout', None)

    def test__load_sets_slidelayout(self):
        """Slide._load() sets slidelayout"""
        # setup -----------------------
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        slidelayout = Mock(name='slideLayout')
        slidelayout.partname = '/ppt/slideLayouts/slideLayout1.xml'
        rel = Mock(name='pptx.packaging._Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT_SLIDE_LAYOUT
        rel.target = slidelayout
        pkgpart = Mock(name='pptx.packaging.Part')
        with open(path, 'rb') as f:
            pkgpart.blob = f.read()
        pkgpart.relationships = [rel]
        part_dict = {slidelayout.partname: slidelayout}
        slide = self.sld._load(pkgpart, part_dict)
        # exercise --------------------
        retval = slide.slidelayout
        # verify ----------------------
        expected = slidelayout
        actual = retval
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test___minimal_element_xml(self):
        """Slide.__minimal_element generates correct XML"""
        # setup -----------------------
        path = os.path.join(thisdir, 'test_files/minimal_slide.xml')
        # exercise --------------------
        sld = self.sld._Slide__minimal_element
        # verify ----------------------
        with open(path, 'r') as f:
            expected_xml = f.read()
        sld_xml = prettify_nsdecls(
            oxml_tostring(sld, encoding='UTF-8', pretty_print=True,
                          standalone=True))
        sld_xml_lines = sld_xml.split('\n')
        expected_xml_lines = expected_xml.split('\n')
        for idx, line in enumerate(sld_xml_lines):
            # msg = '\n\n%s' % sld_xml
            msg = "expected:\n%s\n, got\n%s" % (expected_xml, sld_xml)
            self.assertEqual(line, expected_xml_lines[idx], msg)


class TestSlideCollection(TestCase):
    """Test SlideCollection"""
    def setUp(self):
        prs = Presentation()
        self.slides = SlideCollection(prs)

    def test_add_slide_returns_slide(self):
        """SlideCollection.add_slide() returns instance of Slide"""
        # exercise --------------------
        retval = self.slides.add_slide(None)
        # verify ----------------------
        self.assertIsInstance(retval, Slide)

    def test_add_slide_sets_slidelayout(self):
        """
        SlideCollection.add_slide() sets Slide.slidelayout

        Kind of a throw-away test, but was helpful for initial debugging.
        """
        # setup -----------------------
        slidelayout = Mock(name='slideLayout')
        slidelayout.shapes = []
        slide = self.slides.add_slide(slidelayout)
        # exercise --------------------
        retval = slide.slidelayout
        # verify ----------------------
        expected = slidelayout
        actual = retval
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_add_slide_adds_slide_layout_relationship(self):
        """SlideCollection.add_slide() adds relationship prs->slide"""
        # setup -----------------------
        prs = Presentation()
        slides = prs.slides
        slidelayout = SlideLayout()
        slidelayout._shapes = []
        # exercise --------------------
        slide = slides.add_slide(slidelayout)
        # verify length ---------------
        expected = 1
        actual = len(prs._relationships)
        msg = ("expected len(prs._relationships) of %d, got %d"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)
        # verify values ---------------
        rel = prs._relationships[0]
        expected = ('rId1', RT_SLIDE, slide)
        actual = (rel._rId, rel._reltype, rel._target)
        msg = ("expected relationship 1:, got 2:\n1: %s\n2: %s"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)

    def test_add_slide_sets_partname(self):
        """SlideCollection.add_slide() sets partname of new slide"""
        # setup -----------------------
        prs = Presentation()
        slides = prs.slides
        slidelayout = SlideLayout()
        slidelayout._shapes = []
        # exercise --------------------
        slide = slides.add_slide(slidelayout)
        # verify ----------------------
        expected = '/ppt/slides/slide1.xml'
        actual = slide.partname
        msg = "expected partname '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class TestSlideLayout(TestCase):
    """Test SlideLayout"""
    def setUp(self):
        self.slidelayout = SlideLayout()

    def __loaded_slidelayout(self, prs_slidemaster=None):
        """
        Return SlideLayout instance loaded using mocks. *prs_slidemaster* is
        an already-loaded model-side SlideMaster instance (or mock, as
        appropriate to calling test).
        """
        # partname for related slideMaster
        sldmaster_partname = '/ppt/slideMasters/slideMaster1.xml'
        # path to test slideLayout XML
        slidelayout_path = absjoin(test_file_dir, 'slideLayout1.xml')
        # model-side slideMaster part
        if prs_slidemaster is None:
            prs_slidemaster = Mock(spec=SlideMaster)
        # a part dict containing the already-loaded model-side slideMaster
        loaded_part_dict = {sldmaster_partname: prs_slidemaster}
        # a slideMaster package part for rel target
        pkg_slidemaster_part = Mock(spec=pptx.packaging.Part)
        pkg_slidemaster_part.partname = sldmaster_partname
        # a package-side relationship from slideLayout to its slideMaster
        rel = Mock(name='pptx.packaging._Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT_SLIDE_MASTER
        rel.target = pkg_slidemaster_part
        # the slideLayout package part to send to _load()
        pkg_slidelayout_part = Mock(spec=pptx.packaging.Part)
        pkg_slidelayout_part.relationships = [rel]
        with open(slidelayout_path, 'rb') as f:
            pkg_slidelayout_part.blob = f.read()
        # _load and return
        slidelayout = SlideLayout()
        return slidelayout._load(pkg_slidelayout_part, loaded_part_dict)

    def test__load_sets_slidemaster(self):
        """SlideLayout._load() sets slidemaster"""
        # setup -----------------------
        prs_slidemaster = Mock(spec=SlideMaster)
        # exercise --------------------
        loaded_slidelayout = self.__loaded_slidelayout(prs_slidemaster)
        # verify ----------------------
        expected = prs_slidemaster
        actual = loaded_slidelayout.slidemaster
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_slidemaster_is_readonly(self):
        """SlideLayout.slidemaster is read-only"""
        # verify ----------------------
        self.assertIsReadOnly(self.slidelayout, 'slidemaster')

    def test_slidemaster_raises_on_ref_before_assigned(self):
        """SlideLayout.slidemaster raises on referenced before assigned"""
        with self.assertRaises(AssertionError):
            self.slidelayout.slidemaster


class TestSlideMaster(TestCase):
    """Test SlideMaster"""
    def setUp(self):
        self.sldmaster = SlideMaster()

    def test_slidelayouts_property_empty_on_construction(self):
        """SlideMaster.slidelayouts property empty on construction"""
        # verify ----------------------
        self.assertIsSizedProperty(self.sldmaster, 'slidelayouts', 0)

    def test_slidelayouts_correct_length_after_open(self):
        """SlideMaster.slidelayouts correct length after open"""
        # setup -----------------------
        pkg = Package(test_pptx_path)
        slidemaster = pkg.presentation.slidemasters[0]
        # exercise --------------------
        slidelayouts = slidemaster.slidelayouts
        # verify ----------------------
        self.assertLength(slidelayouts, 11)
